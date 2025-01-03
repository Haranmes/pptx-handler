"""Main module."""
from pptx import Presentation as pp
import pandas as pd
from datetime import datetime
from pathlib import Path
import re
import json
import xlwings as xw
from pptx.util import Pt
import os
import math
from typing import List
import win32com.client as win32



class PowerpointHandler:
    """
    A class to handle PowerPoint presentations, including adding tables, charts, and images.

    Attributes:
        costumer_name (str): The name of the customer.
        powerpoint_dir (Path): The directory where the PowerPoint files are stored.
        powerpoint_imges (Path): The directory where the PowerPoint images are stored.
        logo_path (Path): The path to the logo image.
        pp (Presentation): The PowerPoint presentation object.
        elements (dict): A dictionary to store the elements of each slide.
        elements_file_path (Path): The path to the JSON file storing the elements.
        chart_path_with_names (list): A list of paths to the exported charts.
    """

    def __init__(self, powerpoint_images_dir: list(), costumer_name: str, target_dir: Path):
        """
        Initializes the PowerpointHandler with the given directories and customer name.

        Args:
            powerpoint_images_dir (Path): The directory where the PowerPoint images are stored.
            costumer_name (str): The name of the customer.
            target_dir (Path): The directory where the output PowerPoint files will be saved.
        """
        template_file_name = "202x-xx-xx_Datenanalyse_AKL_Kundenname.pptx"
        self.powerpoint_dir = target_dir
        self.powerpoint_imges = powerpoint_images_dir


        # Get working directory
        self.template_dir = str(Path(__file__).resolve().parent / 'template' / template_file_name)
        self.elements_file_path = self.powerpoint_dir / 'elements.json'

        self.pp = pp(self.template_dir)
        self.elements = self.__get_elements_per_slide()

        self.costumer_name = costumer_name

        current_date = datetime.now().strftime('%Y-%m-%d')
        self.output_path = f"{current_date}_Datenanalyse_AKL_{self.costumer_name}.pptx"
        self.target_dir = target_dir
        self.output_dir = str(self.target_dir / self.output_path)

        self.logo_path = None
        self.pptx_worked_on = False

        for image in self.powerpoint_imges:
            file_path = Path(image)
            if file_path.suffix == '.png':
                file_name = file_path.stem
                if self.like_operator('%ogo%', file_name):
                    self.logo_path = file_path



    def like_operator(self, pattern, string) -> bool:
        """
        Converts SQL LIKE pattern to regex pattern and matches it with the given string.

        Args:
            pattern (str): The SQL LIKE pattern.
            string (str): The string to match.

        Returns:
            bool: True if the string matches the pattern, False otherwise.
        """
        regex_pattern = pattern.replace('%', '.*').replace('_', '.')
        return re.match(regex_pattern, string) is not None

    def __get_elements_per_slide(self) -> dict:
        """
        Extracts the indices of all shapes in the given slide and saves them to a JSON file.

        Returns:
            dict: A dictionary with slide indices as keys and shape indices as values.
        """
        slide_shapes_name = {}
        for slide_idx, slide in enumerate(self.pp.slides):
            slide_shapes_name[slide_idx] = {}
            for shape_idx, shape in enumerate(slide.shapes):
                slide_shapes_name[slide_idx][shape.name] = shape_idx

        with open(self.elements_file_path, 'w') as json_file:
            json.dump(slide_shapes_name, json_file, indent=4)

        print(slide_shapes_name)
        return slide_shapes_name

    def __update_elements_of_slide(self, slide_number: int) -> None:
        """
        Updates the elements of a specific slide and saves the changes to the JSON file.

        Args:
            slide_number (int): The slide number to update.
        """
        slide = self.pp.slides[slide_number] if self.pp.slides[slide_number] is not None else None
        if slide is None:
            raise ValueError("The slide number is not valid.")

        # Add new shapes to self.elements
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.name not in self.elements[slide_number]:
                self.elements[slide_number][shape.name] = shape_idx

        # Remove shapes from self.elements that no longer exist in slide.shapes
        existing_shape_names = {shape.name for shape in slide.shapes}
        for shape_name in list(self.elements[slide_number].keys()):
            if shape_name not in existing_shape_names:
                del self.elements[slide_number][shape_name]

        with open(self.elements_file_path, 'w') as json_file:
            json.dump(self.elements, json_file, indent=4)


    def __get_shape_and_slide(self, slide_number: int, shape_name: str) -> tuple:
        """
        Retrieves the shape and slide for a given slide number and shape name.

        Args:
            slide_number (int): The slide number.
            shape_name (str): The name of the shape.

        Returns:
            tuple: A tuple containing the shape and slide objects.
        """
        slide = self.pp.slides[slide_number] if self.pp.slides[slide_number] is not None else None
        if slide is None:
            raise ValueError("The slide number is not valid.")
        shape_id = self.elements[slide_number].get(shape_name, "Shape not found")
        if shape_name == "Shape not found":
            return None
        shape = slide.shapes[shape_id]
        return shape, slide

    def __separate_row_column(self, cell_reference: str) -> tuple:
        """
        Separates the row and column from a cell reference.

        Args:
            cell_reference (str): The cell reference (e.g., 'C4').

        Returns:
            tuple: A tuple containing the column and row (e.g., ('C', 4)).
        """
        match = re.match(r"([A-Z]+)([0-9]+)", cell_reference, re.I)
        if match:
            column, row = match.groups()
            return column, int(row)
        else:
            raise ValueError("Invalid cell reference")

    def __bring_shape_to_foreground(self, slide_number: int, shape_name: str) -> None:
        """
        Brings the specified shape to the foreground on a slide.

        Args:
            slide_number (int): The slide number where the shape is located.
            shape_name (str): The name of the shape to be brought to the foreground.
        """
        shape, slide = self.__get_shape_and_slide(slide_number, shape_name)
        if shape is None:
            raise ValueError(f"Shape with name {shape_name} not found on slide {slide_number}")

        shape.z_order = 0  # Bring the shape to the front

        self.__update_elements_of_slide(slide_number)
        self.__save_presentation()

    def __save_presentation(self):
        """
        Saves the PowerPoint presentation with the current date and customer name.
        """
        self.pp.save(self.output_dir)
        print(f"Presentation saved to {self.target_dir / self.output_path}")
        self.pptx_worked_on = True
    def add_logo(self, image_path: str, shape_name: str = "logo", slide_number: int = 0, set_to_foreground : bool = False) -> None:
        """
        Adds a logo image to the slide by replacing an existing shape.

        Args:
            image_path (str): The path to the image file to be added.
            shape_name (str): The name of the shape to be replaced.
            slide_number (int): The slide number where the image will be added.
        """
        shape, slide = self.__get_shape_and_slide(slide_number, shape_name)
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        image_path_string = str(image_path)
        slide.shapes.add_picture(image_path_string, left, top, width, height)

        if set_to_foreground:
            self.__bring_shape_to_foreground(slide_number, shape_name)

        self.__update_elements_of_slide(slide_number)

    def add_costumer_name(self, costumer_name: str, slide_number: int = 0, shape_name: str = "costumer", set_to_foreground : bool = False) -> None:
        """
        Adds the customer name to the slide by replacing an existing shape.

        Args:
            costumer_name (str): The name of the customer.
            slide_number (int): The slide number where the name will be added.
            shape_name (str): The name of the shape to be replaced.
        """
        shape, slide = self.__get_shape_and_slide(slide_number, shape_name)
        if shape.has_text_frame is not True:
            raise ValueError("The shape is not a text frame.")
        shape.text = costumer_name

        if set_to_foreground:
            self.__bring_shape_to_foreground(slide_number, shape_name)

        self.__update_elements_of_slide(slide_number)



    def add_table(self, title: str, slide_number: int, table: pd.DataFrame, shape_name: str, set_to_foreground : bool = False) -> None:
        """
        Adds a table to the slide.

        Args:
            title (str): The title of the table.
            slide_number (int): The slide number where the table will be added.
            table (pd.DataFrame): The data to be displayed in the table.
            shape_name (str): The name of the shape to be replaced.
        """
        shape, slide = self.__get_shape_and_slide(slide_number, shape_name)

        left = shape.left
        top = shape.top
        rows, cols = table.shape
        width = shape.width
        height = shape.height
        sp = shape._element
        sp.getparent().remove(sp)

        table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
        for col_idx, col_name in enumerate(table.columns):
            cell = table_shape.cell(0, col_idx)
            cell.text = col_name
        for row_idx, row in table.iterrows():
            for col_idx, value in enumerate(row):
                cell = table_shape.cell(row_idx + 1, col_idx)
                cell.text = str(value)

        if set_to_foreground:
            self.__bring_shape_to_foreground(slide_number, shape_name)

        self.__update_elements_of_slide(slide_number)

    def add_chart_from_excel(self, path_to_excel_file: str, sheet_name: str, slide_number: int, chart_name: str,
                             shape_name: str, set_to_foreground: bool = False) -> None:
        """
        Exports a chart from an Excel file and adds it to a PowerPoint slide.

        Args:
            path_to_excel_file (str): The path to the Excel file.
            sheet_number (int): The sheet number to export charts from.
            slide_number (int): The slide number where the chart will be added.
            shape_name (str): The name of the shape to be replaced.
        """
        # Powerpoint Application works via 1-based indexing instead of 0-based indexing
        slide_number += 1

        current_date = datetime.now().strftime('%Y-%m-%d')
        output_path = f"{current_date}_Datenanalyse_AKL_{self.costumer_name}.pptx"

        xlApp = win32.Dispatch('Excel.Application')
        wb = xlApp.Workbooks.Open(path_to_excel_file)

        pptApp = win32.Dispatch('PowerPoint.Application')
        pptApp.Visible = True

        if self.pptx_worked_on:
            ppt = pptApp.Presentations.Open(self.target_dir / output_path)
        else:
            ppt = pptApp.Presentations.Open(self.template_dir)

        window = pptApp.ActiveWindow
        slide = ppt.Slides.Item(slide_number)
        window.View.GotoSlide(slide_number)

        # Copy the chart from Excel
        wb.Sheets(sheet_name).ChartObjects(chart_name).Copy()

        # Find the shape to replace
        shape_to_replace = None
        for shape in slide.Shapes:
            print(shape.name)
            if shape.name == shape_name:
                shape_to_replace = shape
                print(f"Found shape {shape_name} on slide {slide_number}")
                break

        if shape_to_replace:
            print(f"Replacing shape {shape_name} on slide {slide_number}")
            # Get the position and size of the existing shape
            left = shape_to_replace.Left
            top = shape_to_replace.Top
            width = shape_to_replace.Width
            height = shape_to_replace.Height

            # Delete the existing shape
            shape_to_replace.Delete()

            # Paste the copied chart
            slide.Shapes.Paste()

            # Get the newly pasted shape
            new_shape = slide.Shapes(slide.Shapes.Count)
            new_shape.Left = left
            new_shape.Top = top
            new_shape.Width = width
            new_shape.Height = height

            new_shape.ZOrder(1) # msoSendToBack → https://learn.microsoft.com/en-us/office/vba/api/office.msozordercmd

        wb.Close(SaveChanges=False)
        xlApp.Quit()

        # Release COM objects
        del wb
        del xlApp

        ppt.SaveAs(str(self.target_dir / output_path))

        print(f"Presentation saved to {self.target_dir / output_path}")
        pptApp.Quit()

        # Release COM objects
        del ppt
        del window
        del pptApp

        self.pptx_worked_on = True
        self.__update_elements_of_slide(slide_number)

    def add_table_from_excel(self, slide_number: int, shape_name: str, path_to_excel_file: str, sheet_number: int = 0, set_to_foreground : bool = False) -> None:
        """
        Adds a table to the slide from an Excel file.

        Args:
            slide_number (int): The slide number where the table will be added.
            shape_name (str): The name of the shape to be replaced.
            path_to_excel_file (str): The path to the Excel file.
            sheet_number (int): The sheet number to get the table from.
        """
        shape, slide = self.__get_shape_and_slide(slide_number, shape_name)
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        with xw.App(visible=False) as app:
            book = app.books.open(path_to_excel_file)
            sheet = book.sheets[sheet_number]
            table = sheet.used_range.value
            rows, cols = len(table), len(table[0])
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
            for row_idx, row in enumerate(table):
                for col_idx, value in enumerate(row):
                    cell = table_shape.cell(row_idx, col_idx)
                    cell.text = str(value)
            book.close()

        if set_to_foreground:
            self.__bring_shape_to_foreground(slide_number, shape_name)
        else:
            slide.shapes._spTree.insert(2, table_shape._element)

        self.__update_elements_of_slide(slide_number)



    def add_table_from_excel_range(self, slide_number: int, shape_name: str, path_to_excel_file: str,
                                   sheet_number: int,
                                   start_cell: str, end_cell: str, font_size: float = 9,
                                   skip_header: bool = False,
                                   set_to_foreground: bool = False, is_round : bool = False,
                                   round_columns: List[int] = None) -> None:
        """
        Adds table to the slide from a specified range in an Excel file.

        Args:
            slide_number (int): The slide number where the text will be added.
            shape_name (str): The name of the shape to be replaced.
            path_to_excel_file (str): The path to the Excel file.
            sheet_number (int): The sheet number to get the text from.
            start_cell (str): The starting cell of the range.
            end_cell (str): The ending cell of the range.
            font_size (float): The font size of the text.
        """
        # Open the Excel file and get the specified range
        if self.pptx_worked_on:
            self.pp = pp(self.output_dir)

        shape, slide = self.__get_shape_and_slide(slide_number, shape_name)

        if shape is None:
            raise ValueError(f"Shape with name {shape_name} not found on slide {slide_number}")

        with xw.App(visible=False) as app:
            book = app.books.open(path_to_excel_file)
            sheet = book.sheets[sheet_number]

            # split column from row in the start and end cell
            start_column, start_row = self.__separate_row_column(start_cell)
            end_column, end_row = self.__separate_row_column(end_cell)
            table_data = []

            for row in range(start_row, end_row + 1):
                row_data = []
                for column in range(ord(start_column), ord(end_column) + 1):
                    row_data.append(str(sheet[f"{chr(column)}{row}"].value))
                table_data.append(row_data)

            # Get the position and size of the existing shape
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Remove the existing shape
            sp = shape._element
            sp.getparent().remove(sp)

            # Add a new table with the data from the Excel range


            rows, cols =  (len(table_data) + 1 if skip_header else len(table_data)), len(table_data[0])
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
            for row_idx, row in enumerate(table_data):
                for col_idx, value in enumerate(row):
                    cell = table_shape.cell((row_idx + 1 if skip_header else row_idx), col_idx)
                    if value != "None":
                        if is_round and round_columns is not None and col_idx in round_columns:
                            # convert value to float
                            try:
                                value_float = float(value)

                                # round to two decimal places if value is less than or equal to 1 => percent value
                                if value_float <= 1:
                                    value_round = round(value_float, 2)
                                else:
                                    value_round = round(value_float, 0)
                                    value_round = int(value_round)
                                value = str(value_round)
                            except ValueError:
                                value = str(value)
                        cell.text = value
                        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
                    else:
                        cell.text = ""
                        cell.text_frame.paragraphs[0].font.size = Pt(font_size)

        if set_to_foreground:
            self.__bring_shape_to_foreground(slide_number, shape_name)
        else:
            slide.shapes._spTree.insert(2, shape._element)

        self.__update_elements_of_slide(slide_number)
